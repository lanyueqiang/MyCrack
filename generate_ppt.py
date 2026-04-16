from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)  # 蓝色背景
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_list):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)  # 浅灰背景
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(41, 128, 185)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_table_slide(title, table_data):
    """添加表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(41, 128, 185)
    
    # 表格
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table_shape.cell(i, j)
            cell.text = str(cell_text)
            
            # 格式化
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14) if i == 0 else Pt(12)
                paragraph.font.bold = i == 0
                paragraph.alignment = PP_ALIGN.CENTER
                
                if i == 0:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(41, 128, 185)
                else:
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

# ==================== 开始添加幻灯片 ====================

# 幻灯片1：标题页
add_title_slide(
    "路面裂缝分割网络研究",
    "Network for Robust and High-Accuracy Pavement Crack Segmentation\n\n香港城市大学 | Automation in Construction, 2024"
)

# 幻灯片2：背景与问题
add_content_slide(
    "研究背景与问题",
    [
        "• 路面裂缝及时检测和修复关乎交通安全和道路使用寿命",
        "",
        "• 现有方法的主要问题：",
        "  - 裂缝分割精度不足，鲁棒性差",
        "  - 模型难以准确计算路面状况指数（PCI）",
        "  - 缺乏对未见数据集的泛化能力",
        "",
        "• 传统方法的局限：",
        "  - 手动检测效率低，成本高",
        "  - 传统图像处理方法缺乏有效的特征提取",
        "  - 现有深度学习模型未针对裂缝特性优化"
    ]
)

# 幻灯片3：研究目标
add_content_slide(
    "研究目标",
    [
        "✓ 提高裂缝分割的精度和鲁棒性",
        "",
        "✓ 设计针对裂缝特性的神经网络架构",
        "",
        "✓ 实现跨数据集的泛化性能",
        "",
        "✓ 满足实时应用需求（推理速度 > 30 FPS）",
        "",
        "✓ 自动化路面巡检，降低检测成本"
    ]
)

# 幻灯片4：网络整体架构
add_content_slide(
    "MixCrackNet 网络架构",
    [
        "基于 UNet 框架 + ResNet50 骨干网络",
        "",
        "四大创新核心组件：",
        "  1️⃣ 可变形卷积（Deformable Convolution）",
        "      → 自适应学习裂缝边界",
        "",
        "  2️⃣ 高效多尺度注意机制（EMA）",
        "      → 编码全局信息，强化重要特征",
        "",
        "  3️⃣ Mix Structure",
        "      → 融合浅层位置信息和深层语义信息",
        "",
        "  4️⃣ 加权损失函数",
        "      → 解决样本不均衡问题（权重比 5:1）"
    ]
)

# 幻灯片5：核心创新 - 可变形卷积
add_content_slide(
    "核心创新① 可变形卷积",
    [
        "问题：标准卷积采样固定网格，可能学到无关背景",
        "",
        "解决方案：可变形卷积动态调整采样位置",
        "",
        "优势：",
        "  • 自适应学习裂缝的不规则边界形状",
        "  • 避免采样背景信息",
        "  • 在 Crack500 上提升 +2.36% mIoU",
        "",
        "应用层次：",
        "  • 编码器各层通过可变形卷积提取裂缝特征",
        "  • 解码器采用可变形卷积上采样增强分割精度"
    ]
)

# 幻灯片6：核心创新 - EMA注意机制
add_content_slide(
    "核心创新② 高效多尺度注意机制（EMA）",
    [
        "相比 SENet、CBAM、ECA-Net 的优势：",
        "  ✓ 计算效率更高",
        "  ✓ 参数量更少",
        "  ✓ 通过并联子网络实现跨通道交互",
        "",
        "工作原理：",
        "  • 在多个尺度上计算特征重要性",
        "  • 动态重新加权通道特征",
        "  • 保留全局上下文信息",
        "",
        "性能提升：",
        "  • 显著提升模型表达能力",
        "  • 与其他模块配合效果更优"
    ]
)

# 幻灯片7：核心创新 - Mix Structure
add_content_slide(
    "核心创新③ Mix Structure 设计",
    [
        "问题：单独的浅层或深层特征难以充分表示裂缝",
        "",
        "解决方案：创新性融合浅层和深层特征",
        "",
        "设计细节：",
        "  • 浅层特征：保留精确的位置信息与边界",
        "  • 深层特征：包含丰富的语义信息",
        "  • Mix 模块：多尺度融合优化",
        "",
        "性能指标：",
        "  • mIoU 提升 +2.59%",
        "  • mPA 提升 +5.09%（效果最显著！）",
        "",
        "✨ 这是论文的原创贡献，专门针对裂缝分割设计"
    ]
)

# 幻灯片8：核心创新 - 加权损失函数
add_content_slide(
    "核心创新④ 加权损失函数",
    [
        "背景问题：",
        "  • 路面图像中裂缝像素仅占约 5%",
        "  • 严重的样本不均衡问题",
        "  • 直接训练会导致模型偏向背景类",
        "",
        "解决方案：加权交叉熵损失",
        "  • 权重比：背景 : 裂缝 = 1 : 5",
        "  • 增大裂缝类的学习权重",
        "",
        "效果：",
        "  • Crack500 上提升 +2.74% mIoU",
        "  • 有效缓解类别不均衡，改善分割性能"
    ]
)

# 幻灯片9：实验设置
add_content_slide(
    "实验设置",
    [
        "训练环境：",
        "  • GPU：NVIDIA RTX 4090",
        "  • 框架：PyTorch 2.0.1",
        "  • 优化器：Adam | 学习率：1e-3",
        "  • 训练周期：100 epochs | Batch Size：8",
        "",
        "三个标准数据集：",
        "  1. Crack500：500张高分辨率图像（1500×2000）",
        "  2. CrackForest Dataset (CFD)：118张图像",
        "  3. DeepCrack：527张图像（横纵向、沥青混凝土裂缝）",
        "",
        "对比方法：UNet、PSPNet、HRNet、DeepLab v3+、STDC2-Seg75",
        "",
        "评价指标：mIoU、mPA、Dice 系数"
    ]
)

# 幻灯片10：主要结果对比
result_table = [
    ["数据集", "方法", "mIoU (%)", "Dice (%)", "mPA (%)"],
    ["Crack500", "UNet", "49.85", "45.30", "62.10"],
    ["", "DeepLab v3+", "74.65", "70.25", "81.20"],
    ["", "STDC2-Seg75", "79.25", "75.50", "85.40"],
    ["", "MixCrackNet ✓", "81.58", "79.39", "90.50"],
    ["CrackForest", "STDC2-Seg75", "80.10", "78.20", "88.50"],
    ["", "MixCrackNet ✓", "84.55", "82.15", "93.86"],
    ["DeepCrack", "STDC2-Seg75", "85.20", "83.10", "92.50"],
    ["", "MixCrackNet ✓", "89.01", "87.87", "96.07"],
]

slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "实验结果对比"
for paragraph in title_frame.paragraphs:
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(41, 128, 185)

# 表格
rows, cols = len(result_table), len(result_table[0])
left = Inches(0.3)
top = Inches(1.2)
width = Inches(9.4)
height = Inches(5.8)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

for i, row in enumerate(result_table):
    for j, cell_text in enumerate(row):
        cell = table_shape.cell(i, j)
        cell.text = str(cell_text)
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = i == 0
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(41, 128, 185)
            elif "MixCrackNet" in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(144, 238, 144)  # 浅绿
                paragraph.font.bold = True
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

# 幻灯片11：消融实验
ablation_table = [
    ["模块组合", "mIoU (%)", "Dice (%)", "mPA (%)"],
    ["基线 (UNet)", "77.79", "74.65", "83.90"],
    ["+ 可变形卷积", "80.15", "77.05", "84.50"],
    ["+ 加权损失函数", "80.53", "77.80", "86.64"],
    ["+ Mix Structure", "80.38", "77.80", "88.99"],
    ["全部模块 ✓", "81.58", "79.39", "90.50"],
]

slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "消融实验分析"
for paragraph in title_frame.paragraphs:
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(41, 128, 185)

# 表格
rows, cols = len(ablation_table), len(ablation_table[0])
left = Inches(1.5)
top = Inches(1.3)
width = Inches(7)
height = Inches(5)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

for i, row in enumerate(ablation_table):
    for j, cell_text in enumerate(row):
        cell = table_shape.cell(i, j)
        cell.text = str(cell_text)
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = i == 0
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(41, 128, 185)
            elif "全部模块" in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(144, 238, 144)
                paragraph.font.bold = True
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

# 幻灯片12：泛化性验证
add_content_slide(
    "跨数据集泛化性验证",
    [
        "关键发现：模型具有优异的泛化性能！",
        "",
        "测试场景：",
        "  • 使用 Crack500 训练的模型",
        "  • 在 CrackForest 和 DeepCrack 上测试",
        "",
        "结果表现：",
        "  ✓ CrackForest：mIoU = 84.55%（相对基线提升 5.6%）",
        "  ✓ DeepCrack：mIoU = 89.01%（最高性能）",
        "  ✓ 未见数据集上仍保持优异表现",
        "",
        "意义：",
        "  → 模型不会过度拟合单个数据集",
        "  → 适合实际部署到不同的道路环境"
    ]
)

# 幻灯片13：总结与应用
add_content_slide(
    "核心贡献与应用前景",
    [
        "📌 四大核心贡献：",
        "  1. 高精度鲁棒模型（三个数据集全面最优）",
        "  2. 创新的 Mix Structure 设计（原创方法）",
        "  3. 高效的多尺度注意机制（EMA）",
        "  4. 实时推理性能（48-65 FPS）",
        "",
        "🚀 应用前景：",
        "  • 自动化路面巡检系统",
        "  • 实时路面状况评估（PCI 计算）",
        "  • 降低人工检测成本 80%+",
        "  • 支持多种气候和路面类型",
        "",
        "⭐ 发表期刊：Automation in Construction, Vol. 162, 2024"
    ]
)

# 保存演示文稿
output_path = r"d:\work\research\论文\Crack\路面裂缝分割网络_汇报PPT.pptx"
prs.save(output_path)
print(f"PPT 已成功生成！文件位置：{output_path}")
